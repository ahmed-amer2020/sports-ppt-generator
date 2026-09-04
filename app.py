import streamlit as st
import docx
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import io

# 1. إعدادات الصفحة
st.set_page_config(page_title="مُولد عروض التربية الرياضية الاحترافي", page_icon="🏆", layout="wide")

st.title("🏆 صانع العروض الأكاديمية الاحترافي - تربية رياضية")
st.write("أدخل البيانات التوجيهية للبحث وقم برفع ملف Word للحصول على عرض PowerPoint متكامل مع إكانية المعاينة قبل التحميل.")

# 2. الأسئلة التوجيهية والبيانات
api_key = st.text_input("🔑 أدخل مفتاح Gemini API الخاص بك:", type="password")

st.markdown("---")
st.subheader("📋 البيانات التوجيهية للبحث والعرض")

col1, col2 = st.columns(2)
with col1:
    research_type = st.selectbox(
        "🎯 نوع العرض المطلوب:",
        ["مناقشة رسالة (ماجستير / دكتوراه)", "سيمنار تسجيل خطة بحث (Proposal)", "بحث مرجعي / ورقة عمل"]
    )
    specialty = st.selectbox(
        "⚽ التخصص الرياضي الدقيق:",
        ["التدريب الرياضي وعلوم الحركة", "الإدارة الرياضية والترويح", "المناهج وطرق التدريس", "علوم الصحة والتربية البدنية", "علم النفس الرياضي"]
    )
    slides_count = st.select_slider("📊 عدد الشرائح المطلوب:", options=[5, 8, 10, 12], value=8)

with col2:
    research_title = st.text_input("📌 عنوان البحث الفعلي:", placeholder="اكتب عنوان البحث هنا...")
    researcher_name = st.text_input("👨‍🎓 اسم الباحث واللقب:", placeholder="مثال: الباحث/ أحمد محمود")
    supervisors = st.text_input("👨‍🏫 لجنة الإشراف:", placeholder="مثال: أ.د/ محمد علي، د/ إبراهيم حسن")
    university_info = st.text_input("🏛️ الكلية والجامعة:", placeholder="مثال: كلية التربية الرياضية - جامعة حلوان")

st.subheader("💡 تفاصيل إضافية لدعم جودة الملخص")
col_a, col_b = st.columns(2)
with col_a:
    sample_info = st.text_input("🏋️ عينة البحث والبرنامج التدريبي (إن وجد):", placeholder="مثال: 20 لاعباً - برنامج تدريبي 8 أسابيع")
with col_b:
    main_result = st.text_input("📈 أهم نتيجة إحصائية ترغب بإبرازها:", placeholder="مثال: وجود فروق ذات دلالة إحصائية بنسبة تحسن 15%")

uploaded_file = st.file_uploader("📂 اختر ملف البحث (Docx)", type=["docx"])

def read_docx(file):
    doc = docx.Document(file)
    fullText = []
    for para in doc.paragraphs:
        if para.text.strip():
            fullText.append(para.text.strip())
    return '\n'.join(fullText)

def create_pptx(slides_data, meta_info):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    PRIMARY_COLOR = RGBColor(15, 32, 67)    # أزرق كحلي داكن
    ACCENT_COLOR = RGBColor(212, 175, 55)   # ذهبي فاخر
    TEXT_COLOR = RGBColor(40, 40, 40)        # رمادي داكن للقراءة
    
    for idx, slide_info in enumerate(slides_data):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # ---------------- شريحة العنوان (الغلاف) ----------------
        if idx == 0:
            header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(3.2))
            header_box.fill.solid()
            header_box.fill.fore_color.rgb = PRIMARY_COLOR
            header_box.line.color.rgb = PRIMARY_COLOR
            
            gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(13.333), Inches(0.15))
            gold_line.fill.solid()
            gold_line.fill.fore_color.rgb = ACCENT_COLOR
            gold_line.line.color.rgb = ACCENT_COLOR
            
            tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(2.2))
            tf_title = tx_title.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = meta_info['title'] if meta_info['title'] else slide_info.get("title", "عنوان البحث")
            p_title.font.size = Pt(32)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(255, 255, 255)
            p_title.alignment = PP_ALIGN.RIGHT
            
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.8), Inches(10.333), Inches(3.0))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(245, 247, 250)
            card.line.color.rgb = RGBColor(210, 215, 225)
            
            tx_card = slide.shapes.add_textbox(Inches(1.8), Inches(4.0), Inches(9.733), Inches(2.6))
            tf_card = tx_card.text_frame
            tf_card.word_wrap = True
            
            details = [
                f"إعداد الباحث: {meta_info['researcher']}" if meta_info['researcher'] else "",
                f"تحت إشراف: {meta_info['supervisors']}" if meta_info['supervisors'] else "",
                f"الجهة: {meta_info['university']}" if meta_info['university'] else "",
                f"التخصص: {meta_info['specialty']}"
            ]
            
            for line in details:
                if line:
                    p_card = tf_card.add_paragraph()
                    p_card.text = line
                    p_card.font.size = Pt(20)
                    p_card.font.color.rgb = PRIMARY_COLOR
                    p_card.alignment = PP_ALIGN.RIGHT
                    p_card.space_after = Pt(10)
                    
        # ---------------- باقي الشرائح ----------------
        else:
            top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = PRIMARY_COLOR
            top_bar.line.color.rgb = PRIMARY_COLOR
            
            gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.3), Inches(13.333), Inches(0.08))
            gold_line.fill.solid()
            gold_line.fill.fore_color.rgb = ACCENT_COLOR
            gold_line.line.color.rgb = ACCENT_COLOR
            
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.9))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_info.get("title", "عنوان الشريحة")
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.RIGHT
            
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.2))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            
            content_points = slide_info.get("content", [])
            for i, point in enumerate(content_points):
                p2 = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
                p2.text = f"• {point}"
                p2.font.size = Pt(20)
                p2.font.color.rgb = TEXT_COLOR
                p2.alignment = PP_ALIGN.RIGHT
                p2.space_after = Pt(14)
            
    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

if uploaded_file is not None and api_key:
    if st.button("🚀 إنشاء وتجهيز العرض التقديمي"):
        with st.spinner("جاري قراءة البحث وصياغة الشرائح برؤية أكاديمية..."):
            try:
                text_content = read_docx(uploaded_file)
                genai.configure(api_key=api_key)
                
                available_models = []
                try:
                    for m in genai.list_models():
                        if 'generateContent' in m.supported_generation_methods:
                            available_models.append(m.name)
                except Exception:
                    pass
                
                if not available_models:
                    available_models = ['models/gemini-2.5-flash', 'models/gemini-2.0-flash', 'models/gemini-1.5-flash']
                
                prompt = f"""
                أنت أستاذ ورئيس قسم أكاديمي بكلية التربية الرياضية. قم بتلخيص وإعداد عرض PowerPoint مخصص لـ ({research_type}).
                
                المعطيات والبيانات التوجيهية:
                - التخصص: {specialty}
                - عدد الشرائح المطلوب: {slides_count}
                - عينة البحث/المعالجة: {sample_info if sample_info else "استخرجها من النص"}
                - أهم نتيجة للتركيز عليها: {main_result if main_result else "استخرجها من النص"}
                
                قواعد الإخراج الهامة:
                1. اقسم العرض إلى {slides_count} شرائح بالضبط.
                2. صغ العنوان والمحتوى بأسلوب رصين ومناسب للتقديم أمام لجنة المناقشة.
                3. يجب أن تبدأ كل شريحة بكلمة "شريحة:" يتبعها اسم العنوان المباشر بدون أرقام مفروضة.

                الصيغة المطلوبة:
                شريحة: [عنوان الشريحة]
                - [نقطة جوهرية أولى]
                - [نقطة جوهرية ثانية]

                نص البحث:
                {text_content[:5000]}
                """
                
                response = None
                last_error = ""
                
                for model_name in available_models:
                    try:
                        model = genai.GenerativeModel(model_name)
                        response = model.generate_content(prompt)
                        if response and response.text:
                            break
                    except Exception as err:
                        last_error = str(err)
                        continue
                
                if not response:
                    st.error(f"تعذر توليد النص. الخطأ: {last_error}")
                else:
                    ai_text = response.text
                    
                    slides = []
                    current_slide = None
                    
                    for line in ai_text.split('\n'):
                        line = line.strip()
                        if line.startswith("شريحة:") or line.startswith("الشريحة"):
                            if current_slide and current_slide["content"]:
                                slides.append(current_slide)
                            clean_title = line.split(":", 1)[-1].strip() if ":" in line else line
                            current_slide = {"title": clean_title, "content": []}
                        elif line.startswith("-") and current_slide:
                            point_text = line.lstrip("-").strip()
                            if point_text:
                                current_slide["content"].append(point_text)
                    
                    if current_slide and current_slide["content"]:
                        slides.append(current_slide)
                    
                    meta_info = {
                        'title': research_title,
                        'researcher': researcher_name,
                        'supervisors': supervisors,
                        'university': university_info,
                        'specialty': specialty
                    }
                    
                    # حفظ الشرائح في Session State لمعاينتها
                    st.session_state['generated_slides'] = slides
                    st.session_state['meta_info'] = meta_info
                    st.session_state['pptx_bytes'] = create_pptx(slides, meta_info)
                    
                    st.success("🎉 تم إنشاء العرض التقديمي بنجاح! يمكنك معاينته بالأسفل وتنزيله.")

            except Exception as e:
                st.error(f"حدث خطأ أثناء المعالجة: {e}")

# 3. قسم المعاينة والتحميل
if 'generated_slides' in st.session_state:
    st.markdown("---")
    st.subheader("👁️ معاينة الشرائح قبل التحميل")
    
    slides = st.session_state['generated_slides']
    meta = st.session_state['meta_info']
    
    # تحكم بالتنقل بين الشرائح للمعاينة
    slide_idx = st.slider("اختر الشريحة للمعاينة:", 1, len(slides), 1) - 1
    selected_slide = slides[slide_idx]
    
    # تصميم المعاينة بطريقة مشابهة لشرائح الباوربوينت
    with st.container():
        st.markdown(
            f"""
            <div style="border: 2px solid #0F2043; border-radius: 10px; padding: 20px; background-color: #F8F9FA; direction: rtl; text-align: right;">
                <div style="background-color: #0F2043; color: white; padding: 12px 20px; border-radius: 6px; font-size: 22px; font-weight: bold; margin-bottom: 15px; border-bottom: 4px solid #D4AF37;">
                    {selected_slide['title']}
                </div>
                <div style="font-size: 18px; color: #333; line-height: 1.8; padding: 10px;">
            """,
            unsafe_allow_html=True
        )
        
        if slide_idx == 0 and meta['title']:
            st.markdown(f"**عنوان البحث:** {meta['title']}")
            if meta['researcher']: st.markdown(f"**إعداد الباحث:** {meta['researcher']}")
            if meta['supervisors']: st.markdown(f"**إشراف:** {meta['supervisors']}")
            if meta['university']: st.markdown(f"**الجهة:** {meta['university']}")
        
        for p in selected_slide['content']:
            st.markdown(f"• {p}")
            
        st.markdown("</div></div>", unsafe_allow_html=True)
        
    st.write("")
    # زر التحميل
    st.download_button(
        label="📥 تحميل ملف PowerPoint (.pptx)",
        data=st.session_state['pptx_bytes'],
        file_name="عرض_التربية_الرياضية_الاحترافي.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
